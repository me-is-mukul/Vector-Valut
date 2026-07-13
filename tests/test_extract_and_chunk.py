"""Phase 1 extractors, the cleaner, and page-aware chunking.

The through-line in all of these is **page provenance**. Everything Phase 5 says about
"file X, page 12" is only as true as what happens here, and a citation that points at the
wrong page is worse than none — it looks authoritative and sends the user somewhere wrong.
"""

from __future__ import annotations

from pathlib import Path

import pytest

from osdc.domain.enums import FileType
from osdc.domain.models import ExtractedText, PageSpan
from osdc.pipeline.chunk.chunker import chunk_extracted
from osdc.pipeline.embed.clip_embedder import is_image
from osdc.pipeline.extract.cleaner import clean_extracted, clean_text, find_repeated_lines
from osdc.pipeline.extract.detector import detect_type
from osdc.pipeline.extract.office import DocxExtractor, PptxExtractor
from osdc.pipeline.extract.pdf import PdfExtractor


# --- PDF --------------------------------------------------------------------
def _make_pdf(path: Path, pages: list[str]) -> None:
    import pymupdf

    doc = pymupdf.open()
    for body in pages:
        page = doc.new_page()
        page.insert_text((72, 72), body, fontsize=11)
    doc.save(str(path))
    doc.close()


def test_a_pdf_yields_one_span_per_page(tmp_path: Path) -> None:
    path = tmp_path / "notes.pdf"
    _make_pdf(
        path,
        [
            "Paging is a memory management scheme that eliminates external fragmentation. "
            "The page table maps virtual pages onto physical frames in main memory.",
            "Deadlock arises when four conditions hold together: mutual exclusion, hold and "
            "wait, no preemption, and a circular wait among the processes involved.",
            "Disk scheduling algorithms include first come first served, shortest seek time "
            "first, and the elevator algorithm which sweeps the head across the platter.",
        ],
    )

    extracted = PdfExtractor().extract(path)

    assert [p.page for p in extracted.pages] == [1, 2, 3]
    assert "Paging" in extracted.pages[0].text
    assert "Deadlock" in extracted.pages[1].text
    assert extracted.is_image_based is False


def test_a_text_sparse_pdf_is_treated_as_image_based(tmp_path: Path) -> None:
    """Deliberate, not accidental.

    A PDF averaging only a few characters per page is an image-heavy document — a slide of
    bare titles, a certificate, a diagram with a caption. Its real content is in the
    pictures, so OCR (Phase 2) is the right call and it is correct to flag it.
    """
    path = tmp_path / "titles.pdf"
    _make_pdf(path, ["Intro", "Method", "Results"])

    assert PdfExtractor().extract(path).is_image_based is True


def test_a_pdf_with_no_text_layer_is_flagged_image_based(tmp_path: Path) -> None:
    """A scan. Phase 2's OCR is what rescues it; for now it must be honestly flagged
    rather than classified on an empty string."""
    import pymupdf

    path = tmp_path / "scan.pdf"
    doc = pymupdf.open()
    doc.new_page()  # a blank page — no text layer, exactly like a scanned image
    doc.save(str(path))
    doc.close()

    extracted = PdfExtractor().extract(path)
    assert extracted.is_image_based is True
    assert extracted.char_count == 0


def test_a_corrupt_pdf_does_not_crash_the_pipeline(tmp_path: Path) -> None:
    path = tmp_path / "broken.pdf"
    path.write_bytes(b"%PDF-1.4\nthis is not really a pdf")

    extracted = PdfExtractor().extract(path)
    assert extracted.pages == []  # a Review Queue item, not an exception


def test_pdf_extractor_only_claims_pdfs(tmp_path: Path) -> None:
    extractor = PdfExtractor()
    assert extractor.supports(tmp_path / "a.pdf", FileType.PDF) is True
    assert extractor.supports(tmp_path / "a.docx", FileType.DOCX) is False


# --- DOCX / PPTX ------------------------------------------------------------


def test_docx_extracts_paragraphs_and_tables(tmp_path: Path) -> None:
    import docx

    path = tmp_path / "report.docx"
    document = docx.Document()
    document.add_paragraph("Paging is a memory management scheme.")
    table = document.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Subject"
    table.rows[0].cells[1].text = "Operating Systems"
    document.save(str(path))

    extracted = DocxExtractor().extract(path)

    assert "Paging is a memory management scheme." in extracted.text
    # Tables carry real content in exactly the documents this app is for.
    assert "Operating Systems" in extracted.text
    assert [p.page for p in extracted.pages] == [1], "docx has no readable pagination"


def test_pptx_gives_one_span_per_slide(tmp_path: Path) -> None:
    """Slides DO have real page semantics, so a citation can honestly say 'slide 2'."""
    from pptx import Presentation
    from pptx.util import Inches

    path = tmp_path / "deck.pptx"
    deck = Presentation()
    for title in ("Paging", "Deadlock"):
        slide = deck.slides.add_slide(deck.slide_layouts[5])
        slide.shapes.title.text = title
        box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
        box.text_frame.text = f"Notes about {title.lower()}."
    deck.save(str(path))

    extracted = PptxExtractor().extract(path)

    assert [p.page for p in extracted.pages] == [1, 2]
    assert "Paging" in extracted.pages[0].text
    assert "deadlock" in extracted.pages[1].text.lower()


# --- type detection ---------------------------------------------------------


def test_type_is_detected_from_magic_bytes_not_the_extension(tmp_path: Path) -> None:
    """A scanner emitting `Scan_001.pdf` that is actually a JPEG is common. The extension
    lies; the header does not."""
    lying = tmp_path / "Scan_001.pdf"
    lying.write_bytes(b"\xff\xd8\xff\xe0\x00\x10JFIF\x00" + b"\x00" * 64)

    assert detect_type(lying) is FileType.IMAGE


def test_is_image_uses_the_same_detector(tmp_path: Path) -> None:
    lying = tmp_path / "Scan_002.pdf"
    lying.write_bytes(b"\xff\xd8\xff\xe0\x00\x10JFIF\x00" + b"\x00" * 64)

    assert is_image(lying) is True


def test_text_types_fall_back_to_the_suffix(tmp_path: Path) -> None:
    path = tmp_path / "notes.md"
    path.write_text("# heading", encoding="utf-8")
    assert detect_type(path) is FileType.MD


# --- the cleaner ------------------------------------------------------------


def test_hyphenated_line_breaks_are_rejoined() -> None:
    assert clean_text("inter-\nnational") == "international"


def test_a_real_hyphen_survives() -> None:
    """`Operating-\\nSystems` is a genuine hyphen, not a line-break artefact."""
    assert clean_text("Operating-\nSystems") == "Operating-\nSystems"


def test_running_headers_are_detected_across_pages() -> None:
    pages = [
        PageSpan(page=i, text=f"CS-501 Operating Systems\nContent of page {i}")
        for i in (1, 2, 3, 4)
    ]
    assert "CS-501 Operating Systems" in find_repeated_lines(pages)


def test_running_headers_are_stripped_but_pages_are_kept() -> None:
    """Left in, the header appears in every chunk and every chunk's embedding drifts
    toward it — so all chunks start to look alike and retrieval degrades."""
    extracted = ExtractedText(
        pages=[
            PageSpan(page=i, text=f"CS-501 Operating Systems\nUnique content {i}")
            for i in (1, 2, 3, 4)
        ]
    )
    cleaned = clean_extracted(extracted)

    assert [p.page for p in cleaned.pages] == [1, 2, 3, 4], "page numbers must survive cleaning"
    for page in cleaned.pages:
        assert "CS-501" not in page.text
        assert "Unique content" in page.text


def test_a_two_page_document_keeps_its_headers() -> None:
    """With too few pages, 'repeated' is meaningless and we would strip real content."""
    pages = [PageSpan(page=i, text="Title\nbody") for i in (1, 2)]
    assert find_repeated_lines(pages) == set()


def test_page_number_lines_are_dropped() -> None:
    extracted = ExtractedText(pages=[PageSpan(page=1, text="Real content here.\n12")])
    cleaned = clean_extracted(extracted)
    assert "12" not in cleaned.pages[0].text


# --- chunking ---------------------------------------------------------------


def test_a_chunk_never_crosses_a_page_boundary() -> None:
    """The invariant citations depend on. A chunk spanning two pages cannot honestly
    report which page it came from."""
    extracted = ExtractedText(
        pages=[PageSpan(page=n, text=f"Content for page {n}. " * 30) for n in (1, 2, 3)]
    )
    chunks = chunk_extracted(extracted)

    for chunk in chunks:
        assert f"page {chunk.page}." in chunk.text, (
            f"chunk on page {chunk.page} contains text from another page"
        )


def test_pages_are_split_when_too_long_and_ordinals_stay_sequential() -> None:
    extracted = ExtractedText(pages=[PageSpan(page=1, text="word " * 800)])
    chunks = chunk_extracted(extracted, chunk_chars=400, overlap_chars=50)

    assert len(chunks) > 1
    assert [c.ordinal for c in chunks] == list(range(len(chunks)))
    assert all(c.page == 1 for c in chunks)


def test_a_short_page_stays_a_single_chunk() -> None:
    extracted = ExtractedText(pages=[PageSpan(page=7, text="A short but complete paragraph.")])
    (chunk,) = chunk_extracted(extracted)

    assert chunk.page == 7, "the page number must be carried through, not reset"
    assert chunk.text == "A short but complete paragraph."


def test_a_document_of_only_tiny_pages_falls_back_to_one_chunk_per_page() -> None:
    """Otherwise a slide deck of bare titles indexes to nothing and is unsearchable.

    Critically the fallback emits one chunk *per page*, not one merged chunk for the whole
    document — a merged chunk would have to claim a single page number for content drawn
    from several, and cite the wrong one.
    """
    extracted = ExtractedText(pages=[PageSpan(page=1, text="Intro"), PageSpan(page=2, text="End")])
    chunks = chunk_extracted(extracted)

    assert [(c.ordinal, c.page) for c in chunks] == [(0, 1), (1, 2)]


def test_empty_extraction_produces_no_chunks() -> None:
    assert chunk_extracted(ExtractedText(pages=[])) == []


@pytest.mark.parametrize("size", [50, 200, 1000])
def test_chunks_respect_the_size_budget(size: int) -> None:
    extracted = ExtractedText(pages=[PageSpan(page=1, text="sentence here. " * 200)])
    chunks = chunk_extracted(extracted, chunk_chars=size, overlap_chars=10)
    # Allow the overlap slack, but nothing wildly over budget.
    assert all(len(c.text) <= size * 1.3 for c in chunks)
