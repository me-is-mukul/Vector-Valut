"""DOCX and PPTX extraction.

A note on page numbers, because it affects what citations can honestly say:

- **PPTX has real page semantics.** One slide is one ``PageSpan``, so a citation can say
  "slide 7" and mean it.
- **DOCX does not.** Pagination is a rendering decision made by Word at layout time; the
  file itself has no page breaks to read. So everything lands on page 1 and citations
  fall back to the file plus the chunk's position. Pretending otherwise would mean
  inventing page numbers, and a citation that points at the wrong page is worse than one
  that admits it only knows the file.
"""

from __future__ import annotations

import logging
from pathlib import Path

from osdc.domain.enums import FileType
from osdc.domain.models import ExtractedText, PageSpan

logger = logging.getLogger(__name__)


class DocxExtractor:
    name = "docx"

    def supports(self, path: Path, file_type: FileType) -> bool:
        return file_type is FileType.DOCX

    def extract(self, path: Path) -> ExtractedText:
        import docx

        try:
            document = docx.Document(str(path))
        except Exception as exc:
            logger.warning("python-docx failed on %s: %s", path.name, exc)
            return ExtractedText(pages=[])

        blocks = [p.text.strip() for p in document.paragraphs if p.text.strip()]

        # Tables carry real content in exactly the documents this app is for —
        # mark sheets, fee receipts, lab records — so they are not optional.
        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    blocks.append(" | ".join(cells))

        text = "\n".join(blocks)
        return ExtractedText(
            pages=[PageSpan(page=1, text=text)] if text else [],
            is_image_based=False,
            ocr_used=False,
        )


class PptxExtractor:
    name = "pptx"

    def supports(self, path: Path, file_type: FileType) -> bool:
        return file_type is FileType.PPTX

    def extract(self, path: Path) -> ExtractedText:
        from pptx import Presentation

        try:
            deck = Presentation(str(path))
        except Exception as exc:
            logger.warning("python-pptx failed on %s: %s", path.name, exc)
            return ExtractedText(pages=[])

        pages: list[PageSpan] = []
        for index, slide in enumerate(deck.slides, start=1):
            parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in paragraph.runs).strip()
                        if line:
                            parts.append(line)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))

            # Speaker notes are often where the actual explanation lives.
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes:
                    parts.append(notes)

            if parts:
                pages.append(PageSpan(page=index, text="\n".join(parts)))

        return ExtractedText(pages=pages, is_image_based=False, ocr_used=False)
